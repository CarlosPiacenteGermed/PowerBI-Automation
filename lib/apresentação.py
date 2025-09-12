
from pptx import Presentation
from pptx.util import Inches, Pt
import pandas as pd

# Carregar os dados das tabelas (substitua pelos seus DataFrames reais)
# tabela_final = ...
# tabela_l3m_final = ...

# Abrir apresentação existente
prs = Presentation(f"C:\\Users\\{username}\\Downloads\\RT_Coordenador_Base.pptx")


# Função para inserir tabela em slide
def inserir_tabela_em_slide(slide, df, top=1.5, left=0.5):
    rows, cols = df.shape
    table = slide.shapes.add_table(rows + 1, cols, Inches(left), Inches(top), Inches(9), Inches(0.8)).table

    # Cabeçalhos
    for i, col_name in enumerate(df.columns):
        cell = table.cell(0, i)
        cell.text = str(col_name)
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].font.bold = True

    # Dados
    for i in range(rows):
        for j in range(cols):
            valor = df.iloc[i, j]
            cell = table.cell(i + 1, j)
            cell.text = str(valor)
            cell.text_frame.paragraphs[0].font.size = Pt(11)
from pptx.dml.color import RGBColor

def inserir_tabela_em_slide_com_cores(slide, df, top=1.5, left=0.5):
    rows, cols = df.shape
    table = slide.shapes.add_table(rows + 1, cols, Inches(left), Inches(top), Inches(9), Inches(0.8)).table

    # Definir cores por índice de coluna
    cores = {  
        3: RGBColor(232, 152, 52),
        4: RGBColor(232, 152, 52), 
        5: RGBColor(119, 170, 197), 
        6: RGBColor(119, 170, 197), 
        7: RGBColor(159, 106, 154), 
        8: RGBColor(159, 106, 154),
        9: RGBColor(98, 173, 54), 
        10: RGBColor(98, 173, 54),
    }

    # Cabeçalhos com cor
    for i, col_name in enumerate(df.columns):
        cell = table.cell(0, i)
        cell.text = str(col_name)
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        if i in cores:
            cell.fill.solid()
            cell.fill.fore_color.rgb = cores[i]

    # Dados
    for i in range(rows):
        for j in range(cols):
            valor = df.iloc[i, j]
            cell = table.cell(i + 1, j)
            cell.text = str(valor)
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            
# Inserir tabela principal no slide com título "Resumo Geral"
slide_principal = prs.slides[4]
if slide_principal:
    inserir_tabela_em_slide(slide_principal, tabela_final)

# Inserir tabela L3M no slide com título "Resumo L3M"
slide_l3m = prs.slides[5]
if slide_l3m:
    inserir_tabela_em_slide_com_cores(slide_l3m, tabela_l3m_final)

# Salvar apresentação modificada
prs.save("Apresentacao_Atualizada.pptx")
print("Apresentação atualizada com sucesso!")
